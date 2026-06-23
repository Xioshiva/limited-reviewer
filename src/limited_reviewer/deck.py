"""Build the PowerPoint review deck with python-pptx."""

from __future__ import annotations

import math
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .analyze import GROUP_NAMES, Analysis
from .charts import generate_charts
from .scryfall import ScryfallClient, front_image_url

CARD_ASPECT = 0.717  # width / height of a standard card image

HEADER = RGBColor(0x1B, 0x2A, 0x41)
ACCENT = RGBColor(0x34, 0x55, 0x7A)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xF1, 0xF5)
# Table colours (more contrast than near-white striping).
TBL_STRIPE = RGBColor(0xCF, 0xDA, 0xE8)
TBL_ALT = RGBColor(0xEC, 0xF1, 0xF6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank = self.prs.slide_layouts[6]
        self._dark_ids: set[int] = set()      # slides with a dark full-bleed bg
        self._agenda_slide = None
        self._agenda_runs: list = []

    # -- primitives --------------------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, slide, l, t, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(self, slide, l, t, w, h, lines, size=12, color=HEADER, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        if isinstance(lines, str):
            lines = [lines]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
        return tf

    def _heading(self, slide, title):
        self._rect(slide, 0, 0, SLIDE_W, Inches(0.12), ACCENT)
        self._text(slide, Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.9),
                   title, size=28, bold=True, color=HEADER)

    # -- slide types -------------------------------------------------------
    def title_slide(self, a: Analysis):
        slide = self._slide()
        self._dark_ids.add(id(slide))
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, HEADER)
        self._rect(slide, 0, Inches(4.55), SLIDE_W, Inches(0.06), GOLD)
        self._text(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4),
                   a.set_name, size=46, bold=True, color=WHITE)
        self._text(slide, Inches(0.95), Inches(3.7), Inches(11.5), Inches(0.7),
                   f"Limited Set Review  ·  {a.code}", size=22, color=GOLD)
        self._text(slide, Inches(0.95), Inches(4.8), Inches(11.5), Inches(0.6),
                   f"{a.total} draftable cards   ·   released {a.released_at or 'n/a'}",
                   size=15, color=LIGHT)
        self._text(slide, Inches(0.95), Inches(6.7), Inches(11.5), Inches(0.4),
                   f"Generated {date.today().isoformat()}  ·  data: Scryfall + 17Lands",
                   size=11, color=MUTED)

    def agenda_slide(self, chapters: list[str]):
        slide = self._slide()
        self._agenda_slide = slide
        self._heading(slide, "Contents")
        self._text(slide, Inches(1.1), Inches(1.55), Inches(11), Inches(0.4),
                   "Click a chapter to jump to it.", size=13, color=MUTED)
        tb = slide.shapes.add_textbox(Inches(1.1), Inches(2.05), Inches(11), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True
        self._agenda_runs = []
        for i, title in enumerate(chapters, 1):
            p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
            p.space_after = Pt(16)
            run = p.add_run()
            run.text = f"{i}.   {title}"
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = ACCENT
            self._agenda_runs.append(run)

    def chapter_slide(self, number: int, title: str):
        slide = self._slide()
        self._dark_ids.add(id(slide))
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, HEADER)
        self._rect(slide, Inches(0.9), Inches(3.55), Inches(3.2), Inches(0.06), GOLD)
        self._text(slide, Inches(0.9), Inches(2.5), Inches(11.5), Inches(0.7),
                   f"Chapter {number}", size=22, bold=True, color=GOLD)
        self._text(slide, Inches(0.88), Inches(3.7), Inches(11.5), Inches(1.4),
                   title, size=42, bold=True, color=WHITE)
        return slide

    def link_chapters(self, chapter_slides: list):
        """Turn the agenda lines into internal hyperlinks to the chapter slides."""
        for run, target in zip(self._agenda_runs, chapter_slides):
            if target is None:
                continue
            rId = self._agenda_slide.part.relate_to(target.part, RT.SLIDE)
            rPr = run._r.get_or_add_rPr()
            hlink = rPr.makeelement(qn("a:hlinkClick"), {})
            hlink.set(qn("r:id"), rId)
            hlink.set("action", "ppaction://hlinksldjump")
            rPr.append(hlink)

    def add_page_numbers(self):
        for idx, slide in enumerate(self.prs.slides, 1):
            if idx == 1:
                continue  # no number on the title slide
            color = LIGHT if id(slide) in self._dark_ids else MUTED
            self._text(slide, Inches(12.45), Inches(7.05), Inches(0.7), Inches(0.32),
                       str(idx), size=10, color=color, align=PP_ALIGN.RIGHT)

    def best_commons_uncommons_slide(self, a: Analysis, img):
        slide = self._slide()
        self._heading(slide, "Best commons & uncommons")
        commons = a.top_by_rarity("common", 5)
        uncommons = a.top_by_rarity("uncommon", 5)
        self._text(slide, Inches(0.4), Inches(1.32), Inches(8), Inches(0.3),
                   "Top 5 commons — by GIH win rate", size=14, bold=True, color=ACCENT)
        self._image_row(slide, [(img(c), _wr_caption(a, c)) for c in commons], 1.7, 1.55)
        self._text(slide, Inches(0.4), Inches(4.22), Inches(8), Inches(0.3),
                   "Top 5 uncommons — by GIH win rate", size=14, bold=True, color=ACCENT)
        self._image_row(slide, [(img(c), _wr_caption(a, c)) for c in uncommons], 4.52, 1.55)

    def overview_slide(self, a: Analysis, charts):
        slide = self._slide()
        self._heading(slide, "Overview")
        slide.shapes.add_picture(str(charts["rarity"]), Inches(0.4), Inches(1.5), width=Inches(6.0))
        bullets = [
            f"{a.total} draftable cards", "",
            f"Creatures: {a.creature_count}  ({a.creature_count * 100 // a.total}%)",
            f"Multicolor: {a.multicolor_count}",
        ]
        for label, n in a.type_counts.items():
            bullets.append(f"{label}: {n}")
        self._text(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2),
                   bullets, size=15)

    def color_slide(self, a: Analysis, charts):
        slide = self._slide()
        self._heading(slide, "Color balance")
        slide.shapes.add_picture(str(charts["color"]), Inches(0.4), Inches(1.5), width=Inches(7.6))
        total = a.total
        bullets = ["Cards available in each color:", ""]
        for k in ["W", "U", "B", "R", "G"]:
            n = a.color_counts[k]
            bullets.append(f"{k}: {n}")
        bullets += ["", f"Colorless: {a.color_counts['C']}", f"Multicolor: {a.multicolor_count}"]
        self._text(slide, Inches(8.3), Inches(1.6), Inches(4.6), Inches(5.2), bullets, size=15)

    def curve_slide(self, a: Analysis, charts):
        slide = self._slide()
        self._heading(slide, "Mana curve")
        slide.shapes.add_picture(str(charts["curve"]), Inches(0.4), Inches(1.5), width=Inches(7.6))
        nonland = sum(a.curve.values())
        avg = sum(k * v for k, v in a.curve.items()) / max(nonland, 1)
        cheap = a.curve[1] + a.curve[2] + a.curve[3]
        top = a.curve[6] + a.curve[7]
        bullets = [
            f"Nonland spells: {nonland}", "",
            f"Avg mana value: {avg:.2f}",
            f"1-3 MV: {cheap} ({cheap * 100 // max(nonland,1)}%)",
            f"6+ MV: {top} ({top * 100 // max(nonland,1)}%)",
        ]
        self._text(slide, Inches(8.3), Inches(1.6), Inches(4.6), Inches(5.2), bullets, size=15)

    def functions_slide(self, a: Analysis, charts):
        slide = self._slide()
        self._heading(slide, "Card functions across the set")
        slide.shapes.add_picture(str(charts["categories"]), Inches(0.6), Inches(1.4), width=Inches(8.2))
        self._text(slide, Inches(9.0), Inches(1.7), Inches(4.0), Inches(5.0),
                   ["Counts are based on Scryfall community function tags,",
                    "cross-referenced to this set's draftable cards.", "",
                    "See later slides + card lists to audit each bucket."],
                   size=13, color=MUTED)

    def removal_slide(self, a: Analysis, charts):
        slide = self._slide()
        self._heading(slide, "Removal")
        rd = a.removal_detail
        n = a.category("removal").count
        self._text(slide, Inches(0.5), Inches(1.25), Inches(12), Inches(0.5),
                   f"{n} removal cards  ·  {rd['instant_speed']} castable at instant speed  ·  "
                   f"{rd['unconditional']} unconditional / {rd['conditional']} conditional",
                   size=15, bold=True, color=ACCENT)
        slide.shapes.add_picture(str(charts["removal_color"]), Inches(0.4), Inches(1.9), width=Inches(6.2))
        slide.shapes.add_picture(str(charts["removal_speed"]), Inches(6.9), Inches(1.9), width=Inches(6.0))

    def cardlist_slide(self, title, names, subtitle=None, ncols=3):
        slide = self._slide()
        self._heading(slide, title)
        top = Inches(1.35)
        if subtitle:
            self._text(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
                       subtitle, size=14, bold=True, color=ACCENT)
            top = Inches(1.75)
        if not names:
            self._text(slide, Inches(0.5), top, Inches(12), Inches(1),
                       "No cards found in this category.", size=14, color=MUTED)
            return slide
        per = math.ceil(len(names) / ncols)
        # pick a font size that fits the column height (~5.3in)
        size = max(8, min(13, int(330 / (per * 1.25))))
        col_w = Inches(12.4 / ncols)
        for i in range(ncols):
            chunk = names[i * per:(i + 1) * per]
            if not chunk:
                continue
            left = Inches(0.5 + i * (12.4 / ncols))
            self._text(slide, left, top, col_w, Inches(5.5),
                       [f"•  {nm}" for nm in chunk], size=size)
        return slide

    def image_gallery(self, title, items, cols=7, rows=2):
        """items: list of (image_path|None, caption). Paginates across slides."""
        per = cols * rows
        pages = [items[i:i + per] for i in range(0, len(items), per)] or [[]]
        left0, top0 = 0.35, 1.5
        usable_w, usable_h = 12.63, 5.55
        cell_w = usable_w / cols
        row_h = usable_h / rows
        img_w = min(cell_w - 0.18, (row_h - 0.36) * CARD_ASPECT)
        img_h = img_w / CARD_ASPECT
        for pi, page in enumerate(pages):
            slide = self._slide()
            suffix = f"  ({pi + 1}/{len(pages)})" if len(pages) > 1 else ""
            self._heading(slide, title + suffix)
            for idx, (img, caption) in enumerate(page):
                r, c = divmod(idx, cols)
                x = left0 + c * cell_w + (cell_w - img_w) / 2
                y = top0 + r * row_h
                if img:
                    slide.shapes.add_picture(str(img), Inches(x), Inches(y),
                                             width=Inches(img_w))
                else:
                    self._rect(slide, Inches(x), Inches(y), Inches(img_w), Inches(img_h), LIGHT)
                self._text(slide, Inches(left0 + c * cell_w), Inches(y + img_h + 0.02),
                           Inches(cell_w), Inches(0.3), caption, size=10, bold=True,
                           color=HEADER, align=PP_ALIGN.CENTER)
        return pages

    def _image_row(self, slide, items, top, img_w, gap=0.12, usable_left=0.4, usable_w=12.5):
        img_h = img_w / CARD_ASPECT
        cell_w = img_w + gap
        total_w = max(len(items) * cell_w - gap, 0)
        start = usable_left + (usable_w - total_w) / 2
        for i, (img, caption) in enumerate(items):
            x = start + i * cell_w
            if img:
                slide.shapes.add_picture(str(img), Inches(x), Inches(top), width=Inches(img_w))
            else:
                self._rect(slide, Inches(x), Inches(top), Inches(img_w), Inches(img_h), LIGHT)
            self._text(slide, Inches(x - 0.1), Inches(top + img_h + 0.02),
                       Inches(img_w + 0.2), Inches(0.28), caption, size=9, bold=True,
                       align=PP_ALIGN.CENTER)

    def archetype_winrate_slide(self, a: Analysis, charts):
        if "archetype_winrate" not in charts:
            return False
        slide = self._slide()
        self._heading(slide, "Archetype win rates")
        slide.shapes.add_picture(str(charts["archetype_winrate"]), Inches(0.4), Inches(1.5),
                                 width=Inches(8.4))
        rated = [ar for ar in a.archetypes if ar.get("win_rate") is not None]
        lines = []
        if rated:
            best = max(rated, key=lambda x: x["win_rate"])
            worst = min(rated, key=lambda x: x["win_rate"])
            lines += [f"Best archetype: {best['name']}  ({best['win_rate'] * 100:.1f}%)",
                      f"Weakest: {worst['name']}  ({worst['win_rate'] * 100:.1f}%)", ""]
        for ar in a.color_ratings:
            if ar.is_summary and ar.win_rate is not None:
                lines.append(f"{ar.color_name}: {ar.win_rate * 100:.1f}%")
        self._text(slide, Inches(9.0), Inches(1.7), Inches(4.0), Inches(5.0), lines,
                   size=13, color=MUTED)
        return True

    def archetype_detail_slide(self, a: Analysis, pair: str, img):
        slide = self._slide()
        ar = a.archetype_winrate(pair)
        meta = a.archetype_meta(pair)
        played = f" · {meta['play_rate'] * 100:.1f}% played" if meta and meta.get("play_rate") else ""
        if ar and ar.win_rate is not None:
            head = (f"{ar.color_name} — {ar.win_rate * 100:.1f}% win rate"
                    f"{played} · {ar.games:,} games")
        else:
            head = f"{pair} — win rate n/a"
        self._heading(slide, head)
        commons = a.archetype_cards(pair, "common", 6)
        uncommons = a.archetype_cards(pair, "uncommon", 6)
        self._text(slide, Inches(0.4), Inches(1.32), Inches(8), Inches(0.3),
                   "Key commons — win rate in this archetype", size=14, bold=True, color=ACCENT)
        self._image_row(slide, [(img(c), _arch_wr_caption(a, pair, c)) for c in commons], 1.7, 1.55)
        self._text(slide, Inches(0.4), Inches(4.22), Inches(8), Inches(0.3),
                   "Key uncommons — win rate in this archetype", size=14, bold=True, color=ACCENT)
        self._image_row(slide, [(img(c), _arch_wr_caption(a, pair, c)) for c in uncommons], 4.52, 1.55)

    def evasion_blockers_slide(self, a: Analysis, charts):
        if "evasion_blockers" not in charts:
            return False
        slide = self._slide()
        self._heading(slide, "Evasion & blockers by color")
        slide.shapes.add_picture(str(charts["evasion_blockers"]), Inches(0.5), Inches(1.5),
                                 width=Inches(9.2))
        self._text(slide, Inches(10.0), Inches(1.8), Inches(3.0), Inches(5.0),
                   ["Who gets through and who holds the ground:", "",
                    "Flying / Menace — evasion",
                    "Reach / Deathtouch — defense",
                    "Trample / Vigilance — both", "",
                    "Lopsided evasion often decides which color is the",
                    "aggressor and which must go over the top."],
                   size=13, color=MUTED)
        return True

    def themes_slide(self, a: Analysis):
        from .categories import THEME_KEYS
        rows = [("Theme", "Count", "Top colors")]
        for k in THEME_KEYS:
            cs = a.category(k)
            top = sorted(((v, col) for col, v in cs.by_color.items() if v), reverse=True)[:2]
            rows.append((cs.label, str(cs.count), ", ".join(f"{col} {v}" for v, col in top)))
        slide = self._slide()
        self._heading(slide, "Deck themes (archetype glue)")
        self._text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
                   "Mechanics the set's archetypes are built around.",
                   size=14, color=MUTED)
        self._table(slide, rows, Inches(1.9))
        return True

    def play_draw_slide(self, a: Analysis, charts):
        pd = a.play_draw
        if not pd or "play_draw" not in charts:
            return False
        slide = self._slide()
        self._heading(slide, "On the play vs on the draw")
        slide.shapes.add_picture(str(charts["play_draw"]), Inches(0.5), Inches(1.5),
                                 width=Inches(6.6))
        edge = (pd.win_rate_on_play - 0.5) * 100
        bullets = [
            f"On the play: {pd.win_rate_on_play * 100:.1f}%",
            f"On the draw: {pd.win_rate_on_draw * 100:.1f}%", "",
            f"Play advantage: +{edge:.1f} pts",
            f"Avg game length: {pd.average_game_length:.1f} turns",
            f"Sample: {pd.sample_size:,} games", "",
            "A bigger play advantage points to a faster, more aggressive format;",
            "a smaller one rewards grindy, card-advantage decks.",
        ]
        self._text(slide, Inches(7.5), Inches(1.7), Inches(5.3), Inches(5.0), bullets, size=14)
        return True

    def other_functions_slide(self, a: Analysis):
        slide = self._slide()
        self._heading(slide, "Other functions")
        keys = ["burn", "bounce", "counterspell", "lifegain", "evasion", "tutor",
                "ramp", "combat_trick"]
        rows = [("Function", "Count", "Top colors")]
        for k in keys:
            cs = a.category(k)
            top = sorted(((v, col) for col, v in cs.by_color.items() if v), reverse=True)[:2]
            top_str = ", ".join(f"{col} {v}" for v, col in top)
            rows.append((cs.label, str(cs.count), top_str))
        self._table(slide, rows, Inches(1.6))

    def archetype_slide(self, a: Analysis, charts):
        slide = self._slide()
        self._heading(slide, "Archetypes — color pairs")
        slide.shapes.add_picture(str(charts["archetypes"]), Inches(0.4), Inches(1.5), width=Inches(8.0))
        lines = ["Gold (exactly two-color) cards", "act as archetype signposts.", ""]
        strongest = max(a.archetypes, key=lambda x: x["count"])
        lines.append(f"Most supported: {strongest['pair']} ({strongest['count']})")
        self._text(slide, Inches(8.7), Inches(1.7), Inches(4.3), Inches(5.0), lines,
                   size=14, color=MUTED)

    def methodology_slide(self, a: Analysis):
        slide = self._slide()
        self._heading(slide, "Methodology & caveats")
        bullets = [
            "Cards/images: Scryfall API (Oracle text, colors, rarity, types).",
            "Win rates: 17Lands GIH WR (games-in-hand win rate), 200-game minimum",
            "    before a rate is treated as reliable.",
            "Card pool: unique draftable cards (basic lands + token layouts excluded).",
            "Function tags (removal, combat tricks, ramp, …) come from Scryfall's",
            "    community 'Tagger' oracle tags, cross-referenced to this set.",
            "Removal is split by spell speed and a keyword-based",
            "    conditional/unconditional heuristic, then ranked by win rate.",
            "",
            "Caveat: tags reflect a card's function in general, not its quality in",
            "this specific Limited format. Treat counts as a map, not a tier list —",
            "card images + win rates are shown so you can verify and adjust.",
        ]
        self._text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.2), bullets, size=15)

    def _table(self, slide, rows, top):
        nrows, ncols = len(rows), len(rows[0])
        tbl = slide.shapes.add_table(nrows, ncols, Inches(1.5), top,
                                     Inches(10.3), Inches(0.4 * nrows)).table
        widths = [Inches(5.0), Inches(2.0), Inches(3.3)]
        for i, w in enumerate(widths[:ncols]):
            tbl.columns[i].width = w
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = val
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(13)
                para.font.bold = (r == 0)
                para.font.color.rgb = WHITE if r == 0 else HEADER
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER if r == 0 else (TBL_STRIPE if r % 2 else TBL_ALT)

    def save(self, path: Path):
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))


def _wr_caption(a: Analysis, card: dict) -> str:
    """GIH win rate + pick position (ALSA)."""
    r = a.rating(card)
    if not r or r.gih_wr is None:
        return "no data"
    pick = f" · P{r.alsa:.1f}" if r.alsa is not None else ""
    suffix = "" if r.reliable else " · low n"
    return f"{r.gih_wr * 100:.1f}%{pick}{suffix}"


def _pick_caption(a: Analysis, card: dict) -> str:
    """Pick position (ALSA) + GIH win rate -- the trap/sleeper story."""
    r = a.rating(card)
    alsa = f"{r.alsa:.1f}" if r and r.alsa is not None else "—"
    wr = f"{r.gih_wr * 100:.1f}%" if r and r.gih_wr is not None else "—"
    return f"pick {alsa} · {wr}"


def _arch_wr_caption(a: Analysis, pair: str, card: dict) -> str:
    """In-archetype GIH win rate (+ pick), falling back to the overall rate."""
    r = a.arch_rating(pair, card)
    if not r or r.gih_wr is None:
        return _wr_caption(a, card)
    pick = f" · P{r.alsa:.1f}" if r.alsa is not None else ""
    suffix = "" if r.reliable else " · low n"
    return f"{r.gih_wr * 100:.1f}%{pick}{suffix}"


def _pivot_caption(a: Analysis, card: dict, pairs: list[str]) -> str:
    """How many archetypes a pivot card is good in + its best in-archetype WR."""
    best = a.arch_winrate(pairs[0], card) if pairs else None
    bw = f"{best * 100:.1f}%" if best is not None else "—"
    tag = "/".join(pairs) if len(pairs) <= 2 else f"{len(pairs)} archs"
    return f"{tag} · {bw}"


def build_deck(a: Analysis, output_path: Path, chart_dir: Path,
               client: ScryfallClient) -> Path:
    charts = generate_charts(a, chart_dir)

    def img(card):
        return client.download_image(front_image_url(card))

    chapters = ["The Set at a Glance", "Archetypes", "Removal & Interaction",
                "Draft Signals", "Best Cards"]
    chapter_slides = []

    deck = Deck()
    deck.title_slide(a)
    deck.agenda_slide(chapters)

    # Chapter 1 — The set at a glance.
    chapter_slides.append(deck.chapter_slide(1, chapters[0]))
    deck.overview_slide(a, charts)
    deck.color_slide(a, charts)
    deck.curve_slide(a, charts)
    deck.play_draw_slide(a, charts)

    # Chapter 2 — Archetypes (win rates within each archetype).
    chapter_slides.append(deck.chapter_slide(2, chapters[1]))
    deck.archetype_winrate_slide(a, charts)
    deck.archetype_slide(a, charts)
    pairs = sorted(a.archetypes,
                   key=lambda x: (x.get("win_rate") is None, -(x.get("win_rate") or 0)))
    for ar in pairs:
        deck.archetype_detail_slide(a, ar["pair"], img)
    for rarity, label in (("common", "commons"), ("uncommon", "uncommons")):
        pivots = a.pivot_cards(rarity, n=14)
        if pivots:
            deck.image_gallery(
                f"Pivot {label} — strong in 2+ archetypes  (archetypes · best WR)",
                [(img(c), _pivot_caption(a, c, prs)) for c, prs in pivots])

    # Chapter 3 — Removal & interaction (image galleries by color, sorted by WR).
    chapter_slides.append(deck.chapter_slide(3, chapters[2]))
    for group, cards in a.grouped_by_color(a.category("removal").cards):
        items = [(img(c), _wr_caption(a, c)) for c in cards]
        deck.image_gallery(f"Removal — {GROUP_NAMES[group]} ({len(cards)}) · by GIH win rate",
                           items)
    sweepers = a.sorted_by_winrate(a.category("board_wipe").cards)
    if sweepers:
        deck.image_gallery(f"Sweepers / Wraths ({len(sweepers)}) · by GIH win rate",
                           [(img(c), _wr_caption(a, c)) for c in sweepers])
    for group, cards in a.grouped_by_color(a.category("combat_trick").cards):
        items = [(img(c), _wr_caption(a, c)) for c in cards]
        deck.image_gallery(f"Combat tricks — {GROUP_NAMES[group]} ({len(cards)}) · by GIH win rate",
                           items)
    deck.other_functions_slide(a)

    # Chapter 4 — Draft signals (traps / sleepers).
    chapter_slides.append(deck.chapter_slide(4, chapters[3]))
    traps, sleepers = a.traps_and_sleepers(8)
    if traps:
        deck.image_gallery(
            "Overrated — taken early, underperform  (pick = ALSA · GIH win rate)",
            [(img(c), _pick_caption(a, c)) for c in traps], cols=8, rows=1)
    if sleepers:
        deck.image_gallery(
            "Underrated sleepers — wheel, but overperform  (pick = ALSA · GIH win rate)",
            [(img(c), _pick_caption(a, c)) for c in sleepers], cols=8, rows=1)

    # Chapter 5 — Best cards.
    chapter_slides.append(deck.chapter_slide(5, chapters[4]))
    deck.best_commons_uncommons_slide(a, img)
    top = a.top_by_winrate(20)
    if top:
        items = [(img(c), f"#{i}  {_wr_caption(a, c)}") for i, c in enumerate(top, 1)]
        deck.image_gallery("Top 20 cards in the set — by GIH win rate", items,
                           cols=10, rows=2)

    deck.methodology_slide(a)
    deck.link_chapters(chapter_slides)
    deck.add_page_numbers()
    deck.save(output_path)
    return output_path
